"""
Pardus AI GUI - Flask backend server.
Serves the Web UI and provides REST API endpoints for AI communication.
"""
import os
import threading
import time
from flask import Flask, render_template, request, jsonify, send_from_directory
from src.core.qwen_api import QwenAPI
from src.core.puter_api import PuterAPI
from src.features.screen_capture import ScreenCapture

# Flask app setup
template_dir = os.path.join(os.path.dirname(__file__), 'templates')
static_dir = os.path.join(os.path.dirname(__file__), 'static')
app = Flask(__name__, template_folder=template_dir, static_folder=static_dir)
app.config['SEND_FILE_MAX_AGE_DEFAULT'] = 0

# AI Services - Qwen Free API üzerinden
_qwen_api = None
_puter_api = None
_screen_capturer = None
_computer_agent = None

def _get_api():
    global _qwen_api
    if _qwen_api is None:
        _qwen_api = QwenAPI()
    return _qwen_api

def _get_puter_api():
    global _puter_api
    if _puter_api is None:
        _puter_api = PuterAPI()
    return _puter_api

def _get_pollinations():
    """Geriye uyumluluk — artık QwenAPI kullanıyor."""
    return _get_api()

def _get_vision_api():
    """Vision API artık Qwen VL modelleri kullanıyor."""
    return _get_api()

def _get_screen_capturer():
    global _screen_capturer
    if _screen_capturer is None:
        _screen_capturer = ScreenCapture()
    return _screen_capturer

import json

def get_chats_file_path():
    config_dir = os.path.expanduser("~/.config/pardus_ai")
    os.makedirs(config_dir, exist_ok=True)
    return os.path.join(config_dir, "chats.json")

# ─── Puter Token Setup ──────────────────────────────────

@app.route('/setup')
def qwen_setup_page():
    """Qwen API durum sayfası."""
    return '''<!DOCTYPE html>
<html><head>
<meta charset="UTF-8">
<title>Pardus AI - Qwen API Durumu</title>
<style>
body { font-family: 'Segoe UI', sans-serif; background: linear-gradient(135deg, #0d1117, #161b22); color: #e6edf3; display: flex; justify-content: center; align-items: center; min-height: 100vh; margin: 0; }
.container { text-align: center; max-width: 520px; padding: 40px; }
.logo { font-size: 64px; margin-bottom: 20px; }
h1 { color: #c084fc; font-size: 24px; margin-bottom: 10px; }
p { color: #8b949e; line-height: 1.6; margin-bottom: 15px; }
.status-box { background: rgba(255,255,255,0.05); border-radius: 12px; padding: 20px 25px; margin: 20px 0; border: 1px solid rgba(255,255,255,0.1); }
.success { color: #3fb950; font-size: 18px; }
.error { color: #f85149; font-size: 18px; }
.btn { padding: 12px 30px; background: linear-gradient(135deg, rgba(139,92,246,0.4), rgba(6,182,212,0.3)); border: 1px solid rgba(139,92,246,0.5); border-radius: 10px; color: #e0e0e0; font-size: 15px; font-weight: 600; cursor: pointer; transition: all 0.2s; }
.btn:hover { background: linear-gradient(135deg, rgba(139,92,246,0.6), rgba(6,182,212,0.45)); transform: translateY(-1px); }
#result { margin-top: 15px; min-height: 24px; }
</style>
</head><body>
<div class="container">
    <div class="logo">✨</div>
    <h1>Pardus AI — Qwen API Durumu</h1>
    <p>Qwen Free API sunucusu <b>port 3264</b> üzerinde çalışıyor.</p>
    <div class="status-box">
        <p id="status">Kontrol ediliyor...</p>
    </div>
    <button class="btn" onclick="checkStatus()">Durumu Kontrol Et</button>
</div>
<script>
async function checkStatus() {
    try {
        const r = await fetch('/api/qwen/status');
        const data = await r.json();
        if (data.available) {
            document.getElementById('status').innerHTML = '<span class="success">✅ Qwen API aktif! ' + (data.model_count || 0) + ' model mevcut.</span>';
        } else {
            document.getElementById('status').innerHTML = '<span class="error">❌ Qwen API sunucusuna ulaşılamıyor.</span>';
        }
    } catch(e) {
        document.getElementById('status').innerHTML = '<span class="error">❌ Bağlantı hatası: ' + e.message + '</span>';
    }
}
checkStatus();
</script>
</body></html>'''

@app.route('/api/qwen/status', methods=['GET'])
def qwen_status():
    """Qwen API durumu."""
    from src.core.qwen_api import QwenAPI
    available = QwenAPI.is_available()
    models = QwenAPI.get_models() if available else []
    return jsonify({
        'available': available,
        'model_count': len(models),
        'models': models[:10]  # İlk 10 modeli göster
    })

# ─── Routes ─────────────────────────────────────────────

@app.route('/api/chats', methods=['GET'])
def get_chats():
    """Load chat history from disk."""
    chats_file = get_chats_file_path()
    if os.path.exists(chats_file):
        try:
            with open(chats_file, 'r', encoding='utf-8') as f:
                return jsonify({'success': True, 'chats': json.load(f)})
        except Exception as e:
            return jsonify({'success': False, 'error': str(e)})
    return jsonify({'success': True, 'chats': {}})

@app.route('/api/chats', methods=['POST'])
def save_chats():
    """Save chat history to disk."""
    try:
        data = request.get_json()
        chats = data.get('chats', {})
        chats_file = get_chats_file_path()
        with open(chats_file, 'w', encoding='utf-8') as f:
            json.dump(chats, f, ensure_ascii=False, indent=2)
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/')
def index():
    return render_template('index.html')

import subprocess
import re

# System prompt for command execution capability
SYSTEM_PROMPT = (
    "Sen 'Pardus AI Asistanı' adlı bir yapay zeka asistanısın. "
    "Pardus Linux için özel olarak geliştirildin. "
    "Kendini asla Claude, GPT, Gemini veya başka bir AI olarak tanıtma. "
    "Sen sadece 'Pardus AI Asistanı'sın.\n\n"
    "Kullanıcının sistem hakkında sorular sorduğunda veya bir görev istediğinde, "
    "gerekli Linux komutlarını çalıştırabilirsin.\n\n"
    "Bir komut çalıştırman gerekiyorsa, yanıtında komutu şu formatta yaz:\n"
    "<cmd>komut buraya</cmd>\n\n"
    "Kurallar:\n"
    "- Birden fazla komut gerekiyorsa her birini ayrı <cmd> etiketinde yaz.\n"
    "- Tehlikeli komutlar (rm -rf /, mkfs, dd if=/dev/zero vb.) ASLA kullanma.\n"
    "- Komut gerektirmeyen normal sohbet sorularında <cmd> etiketi KULLANMA.\n"
    "- Kullanıcının yazdığı dilde yanıt ver.\n"
    "- Kısa ve öz ol."
)

# Dangerous command patterns
BLOCKED_COMMANDS = [
    'rm -rf /', 'rm -rf /*', 'mkfs', ':(){ :|:& };:', 'dd if=/dev/zero',
    'dd if=/dev/random', '> /dev/sda', 'chmod -R 777 /', 'shutdown', 'reboot',
    'init 0', 'init 6', 'halt', 'poweroff', 'mv / ', 'wget', 'curl.*|.*bash',
    'fork bomb'
]

def _is_command_safe(cmd):
    """Check if a command is safe to execute."""
    cmd_lower = cmd.lower().strip()
    for blocked in BLOCKED_COMMANDS:
        if blocked in cmd_lower:
            return False
    return True

def _execute_command(cmd, timeout=15):
    """Execute a shell command safely and return the output."""
    try:
        result = subprocess.run(
            cmd,
            shell=True,
            capture_output=True,
            text=True,
            timeout=timeout,
            executable='/bin/bash',
            cwd=os.path.expanduser('~')
        )
        output = result.stdout
        if result.stderr:
            output += ('\n' if output else '') + result.stderr
        output = output.strip()
        # Truncate very long outputs
        if len(output) > 3000:
            output = output[:3000] + '\n... (çıktı kısaltıldı)'
        return output if output else '(Komut çalıştı, çıktı yok)'
    except subprocess.TimeoutExpired:
        return '(Komut zaman aşımına uğradı)'
    except Exception as e:
        return f'(Hata: {str(e)})'

@app.route('/api/chat', methods=['POST'])
def chat():
    """Handle text chat messages with optional command execution."""
    data = request.get_json()
    user_message = data.get('message', '').strip()
    needs_title = data.get('needs_title', False)
    model_choice = data.get('model', 'cloud')
    
    if not user_message:
        return jsonify({'success': False, 'error': 'Boş mesaj gönderilemez.'})
    
    try:
        # Choose AI backend based on user preference
        full_prompt = f"{SYSTEM_PROMPT}\n\nKullanıcı: {user_message}"
        
        if model_choice == 'local':
            # Local AI mode – use ModelLoader
            try:
                from src.core.model_loader import ModelLoader
                loader = ModelLoader()
                ai_response = loader.generate_response(full_prompt)
            except Exception as local_err:
                return jsonify({'success': False, 'error': f'Yerel model hatası: {str(local_err)}. Ayarlardan Bulut AI moduna geçmeyi deneyin.'})
        else:
            # Cloud AI mode (default)
            messages = [{"role": "user", "content": [{"type": "text", "text": full_prompt}]}]
            ai_response = _get_pollinations().generate_response(messages)
        
        # Step 2: Check if AI requested command execution
        cmd_pattern = re.findall(r'<cmd>(.*?)</cmd>', ai_response, re.DOTALL)
        
        if cmd_pattern:
            # Execute commands and collect outputs
            cmd_outputs = []
            for cmd in cmd_pattern:
                cmd = cmd.strip()
                if _is_command_safe(cmd):
                    output = _execute_command(cmd)
                    cmd_outputs.append(f"$ {cmd}\n{output}")
                else:
                    cmd_outputs.append(f"$ {cmd}\n⚠️ Güvenlik nedeniyle bu komut engellendi.")
            
            all_outputs = '\n\n'.join(cmd_outputs)
            
            # Step 3: Feed command outputs back to AI for interpretation
            interpret_prompt = (
                f"Kullanıcı şunu sormuştu: {user_message}\n\n"
                f"Aşağıdaki komut(lar) çalıştırıldı ve çıktıları şöyle:\n\n{all_outputs}\n\n"
                f"Bu çıktılara dayanarak kullanıcıya Türkçe, anlaşılır ve detaylı bir yanıt ver. "
                f"Komut çıktılarını da göster."
            )
            interpret_messages = [{"role": "user", "content": [{"type": "text", "text": interpret_prompt}]}]
            final_response = _get_pollinations().generate_response(interpret_messages)
        else:
            # No commands needed, use response directly
            final_response = ai_response
        
        result = {'success': True, 'response': final_response}
        
        # Generate title in the same call if requested
        if needs_title:
            try:
                title_prompt = (
                    f"Aşağıdaki sohbetin konusunu özetleyen çok kısa bir başlık yaz (en fazla 5 kelime, "
                    f"sadece başlığı yaz, başka açıklama yapma, tırnak işareti kullanma):\n"
                    f"Kullanıcı: {user_message[:200]}\n"
                    f"Asistan: {final_response[:200]}"
                )
                title_messages = [{"role": "user", "content": [{"type": "text", "text": title_prompt}]}]
                title = _get_pollinations().generate_response(title_messages)
                title = title.strip().strip('"').strip("'").strip()
                if len(title) > 50:
                    title = title[:47] + '...'
                result['title'] = title
            except:
                result['title'] = user_message[:35] + ('...' if len(user_message) > 35 else '')
        
        return jsonify(result)
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/api/screenshot', methods=['POST'])
def screenshot():
    """Take a screenshot and analyze it."""
    try:
        # Small delay so the window can minimize if needed
        time.sleep(0.3)
        
        capturer = _get_screen_capturer()
        project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__))))
        screenshot_path = os.path.join(project_root, 'current_screen.png')
        result_path = capturer.capture_full_screen(screenshot_path)
        
        if not result_path:
            return jsonify({'success': False, 'error': 'Ekran görüntüsü alınamadı.'})
        
        abs_path = os.path.abspath(result_path)
        prompt = "Lütfen bu ekran görüntüsünü detaylıca analiz et ve ne gördüğünü Türkçe anlat."
        
        response = _get_vision_api().generate_vision_response(prompt, abs_path)
        return jsonify({'success': True, 'response': response, 'image_path': abs_path})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/api/analyze-image', methods=['POST'])
def analyze_image():
    """Analyze an uploaded image."""
    try:
        if 'image' not in request.files:
            return jsonify({'success': False, 'error': 'Görsel yüklenmedi.'})
        
        file = request.files['image']
        prompt = request.form.get('prompt', 'Bu görseli açıkla.')
        
        # Save uploaded file temporarily
        project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__))))
        upload_path = os.path.join(project_root, 'uploaded_image.png')
        file.save(upload_path)
        
        response = _get_vision_api().generate_vision_response(prompt, upload_path)
        return jsonify({'success': True, 'response': response})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/api/analyze-video', methods=['POST'])
def analyze_video():
    """Analyze an uploaded video by extracting frames."""
    try:
        if 'video' not in request.files:
            return jsonify({'success': False, 'error': 'Video yüklenmedi.'})
        
        file = request.files['video']
        prompt = request.form.get('prompt', 'Bu videoyu analiz et ve açıkla.')
        
        # Save uploaded video temporarily
        project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__))))
        ext = os.path.splitext(file.filename)[1] if file.filename else '.mp4'
        raw_path = os.path.join(project_root, f'uploaded_video_raw{ext}')
        upload_path = os.path.join(project_root, 'uploaded_video.mp4')
        file.save(raw_path)
        
        print(f"[Video] Kaydedildi: {raw_path} ({os.path.getsize(raw_path)} bytes)")
        
        # Re-encode with ffmpeg to ensure OpenCV compatibility
        import subprocess
        reencode = subprocess.run(
            ['ffmpeg', '-y', '-i', raw_path, '-c:v', 'libx264', '-preset', 'fast',
             '-c:a', 'aac', '-movflags', '+faststart', upload_path],
            capture_output=True, text=True, timeout=300
        )
        
        # Clean up raw file
        if os.path.exists(raw_path):
            os.remove(raw_path)
        
        if reencode.returncode != 0 or not os.path.exists(upload_path):
            print(f"[Video] ffmpeg hatası: {reencode.stderr[:200] if reencode.stderr else 'bilinmiyor'}")
            return jsonify({'success': False, 'error': 'Video dönüştürülemedi. Lütfen farklı bir format deneyin.'})
        
        print(f"[Video] Dönüştürüldü: {upload_path} ({os.path.getsize(upload_path)} bytes)")
        
        try:
            response = _get_vision_api().generate_video_response(prompt, upload_path)
        finally:
            # Clean up temp video file
            if os.path.exists(upload_path):
                os.remove(upload_path)
        
        return jsonify({'success': True, 'response': response})
    except Exception as e:
        print(f"[Video] Hata: {str(e)}")
        return jsonify({'success': False, 'error': str(e)})

@app.route('/api/analyze-audio', methods=['POST'])
def analyze_audio_file():
    """Analyze an uploaded audio file."""
    try:
        if 'audio' not in request.files:
            return jsonify({'success': False, 'error': 'Ses dosyası yüklenmedi.'})
        
        file = request.files['audio']
        prompt = request.form.get('prompt', 'Bu ses dosyasını analiz et ve açıkla.')
        
        # Save uploaded audio temporarily
        project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(__file__))))
        ext = os.path.splitext(file.filename)[1] if file.filename else '.mp3'
        raw_path = os.path.join(project_root, f'uploaded_audio_raw{ext}')
        wav_path = os.path.join(project_root, 'uploaded_audio.wav')
        file.save(raw_path)
        
        print(f"[Audio] Kaydedildi: {raw_path} ({os.path.getsize(raw_path)} bytes)")
        
        import subprocess as sp
        
        # Convert to WAV for analysis
        sp.run(
            ['ffmpeg', '-y', '-i', raw_path, '-acodec', 'pcm_s16le',
             '-ar', '16000', '-ac', '1', wav_path],
            capture_output=True, text=True, timeout=120
        )
        
        if os.path.exists(raw_path):
            os.remove(raw_path)
        
        if not os.path.exists(wav_path) or os.path.getsize(wav_path) < 200:
            return jsonify({'success': False, 'error': 'Ses dosyası dönüştürülemedi.'})
        
        # Speech transcription
        transcript = ""
        try:
            import speech_recognition as sr
            from pydub import AudioSegment
            
            audio_seg = AudioSegment.from_wav(wav_path)
            recognizer = sr.Recognizer()
            full_text = []
            chunk_ms = 30000
            
            for i in range(0, len(audio_seg), chunk_ms):
                chunk = audio_seg[i:i + chunk_ms]
                chunk_path = wav_path + f'_chunk_{i}.wav'
                chunk.export(chunk_path, format='wav')
                try:
                    with sr.AudioFile(chunk_path) as source:
                        audio_data = recognizer.record(source)
                        try:
                            text = recognizer.recognize_google(audio_data, language='tr-TR')
                        except sr.UnknownValueError:
                            try:
                                text = recognizer.recognize_google(audio_data, language='en-US')
                            except sr.UnknownValueError:
                                text = ""
                        if text:
                            full_text.append(text)
                except Exception:
                    pass
                finally:
                    if os.path.exists(chunk_path):
                        os.remove(chunk_path)
            
            transcript = ' '.join(full_text)
            if transcript:
                print(f"[Audio] Transkripsiyon: {len(transcript)} karakter")
        except Exception as e:
            print(f"[Audio] Transkripsiyon hatası: {str(e)[:80]}")
        
        # Sound classification
        audio_desc = ""
        try:
            from src.core.audio_analyzer import analyze_audio
            result = analyze_audio(wav_path)
            if result['has_audio'] and result['description']:
                audio_desc = result['description']
                print(f"[Audio] Sınıflandırma: {audio_desc[:80]}...")
        except Exception as e:
            print(f"[Audio] Sınıflandırma hatası: {str(e)[:80]}")
        
        # Clean up wav
        if os.path.exists(wav_path):
            os.remove(wav_path)
        
        # Generate spectrogram image for visual analysis
        spectrogram_path = os.path.join(project_root, 'audio_spectrogram.png')
        has_spectrogram = False
        try:
            from src.core.audio_analyzer import generate_spectrogram
            # Re-convert for spectrogram (need the file again)
            raw_for_spec = os.path.join(project_root, f'uploaded_audio_raw{ext}')
            if not os.path.exists(raw_for_spec):
                # File was already deleted, use the original upload
                file.seek(0)
                file.save(raw_for_spec)
            has_spectrogram = generate_spectrogram(raw_for_spec, spectrogram_path)
            if os.path.exists(raw_for_spec):
                os.remove(raw_for_spec)
            if has_spectrogram:
                print(f"[Audio] Spektrogram oluşturuldu: {spectrogram_path}")
        except Exception as e:
            print(f"[Audio] Spektrogram hatası: {str(e)[:80]}")
        
        # Build AI prompt — natural, conversational
        ai_prompt_parts = ["Kullanıcı bir ses dosyası yükledi ve analiz etmeni istiyor."]
        
        if transcript:
            ai_prompt_parts.append(f"\nSesteki konuşma:\n\"{transcript}\"")
        
        if audio_desc:
            import re
            duration_match = re.search(r'Ses süresi: ([\d.]+)', audio_desc)
            if duration_match:
                ai_prompt_parts.append(f"\nSes süresi: {duration_match.group(1)} saniye.")
        
        if has_spectrogram:
            ai_prompt_parts.append("\nAşağıda sesin mel spektrogram görüntüsü var. Buna bakarak sesin ne olduğunu doğal bir şekilde açıkla.")
        
        if not transcript and not audio_desc and not has_spectrogram:
            ai_prompt_parts.append("\nSes dosyasından bilgi çıkarılamadı.")
        
        ai_prompt_parts.append(f"\nKullanıcının sorusu: {prompt}")
        
        full_prompt = ' '.join(ai_prompt_parts)
        
        # If we have a spectrogram, send to vision API; otherwise text API
        if has_spectrogram:
            response = _get_vision_api().generate_vision_response(full_prompt, spectrogram_path)
            if os.path.exists(spectrogram_path):
                os.remove(spectrogram_path)
        else:
            messages = [{"role": "user", "content": [{"type": "text", "text": full_prompt}]}]
            response = _get_pollinations().generate_response(messages)
        
        return jsonify({'success': True, 'response': response})
    except Exception as e:
        print(f"[Audio] Hata: {str(e)}")
        return jsonify({'success': False, 'error': str(e)})

# ═══════════════ AGENT MODE API ═══════════════

def _get_agent():
    global _computer_agent
    if _computer_agent is None:
        from src.features.computer_agent import ComputerAgent
        _computer_agent = ComputerAgent(_get_vision_api())
    return _computer_agent

@app.route('/api/agent/start', methods=['POST'])
def agent_start():
    """Start the computer agent with a task."""
    data = request.get_json()
    task = data.get('task', '') if data else ''
    if not task:
        return jsonify({'success': False, 'error': 'Görev açıklaması gerekli.'})
    result = _get_agent().start(task)
    if 'error' in result:
        return jsonify({'success': False, 'error': result['error']})
    return jsonify({'success': True, 'message': result['message']})

@app.route('/api/agent/stop', methods=['POST'])
def agent_stop():
    """Stop the running agent."""
    result = _get_agent().stop()
    return jsonify({'success': True, 'message': result['message']})

@app.route('/api/agent/respond', methods=['POST'])
def agent_respond():
    """Send user response to agent (for password/input requests)."""
    data = request.get_json()
    response_text = data.get('response', '') if data else ''
    result = _get_agent().respond(response_text)
    if 'error' in result:
        return jsonify({'success': False, 'error': result['error']})
    return jsonify({'success': True})

@app.route('/api/agent/status', methods=['GET'])
def agent_status():
    """Get current agent status."""
    status = _get_agent().status
    return jsonify(status)

def _extract_file_text(file_path, ext):
    """Extract text content from uploaded files."""
    ext = ext.lower()
    
    if ext == '.pdf':
        try:
            import PyPDF2
            text = ''
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + '\n'
            return text.strip() if text.strip() else '[PDF dosyasından metin çıkarılamadı]'
        except Exception as e:
            return f"[PDF okuma hatası: {e}]"
    
    elif ext == '.docx':
        try:
            import docx
            doc = docx.Document(file_path)
            return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            return f"[DOCX okuma hatası: {e}]"
    
    elif ext == '.doc':
        try:
            import re
            # Use 'strings' command to extract text from binary .doc file
            # -e l = little-endian 16-bit (how Word stores Unicode text)
            result = subprocess.run(
                ['strings', '-e', 'l', file_path],
                capture_output=True, text=True, timeout=10
            )
            text = result.stdout.strip()
            
            if not text:
                # Fallback: try standard ASCII strings
                result = subprocess.run(
                    ['strings', file_path],
                    capture_output=True, text=True, timeout=10
                )
                text = result.stdout.strip()
            
            if text:
                # Filter out short lines (metadata noise) and keep actual content
                lines = text.split('\n')
                content_lines = [l.strip() for l in lines if len(l.strip()) > 2 
                                 and not re.match(r'^[\x00-\x1f\s]+$', l)]
                return '\n'.join(content_lines) if content_lines else '[DOC dosyasından metin çıkarılamadı]'
            return '[DOC dosyasından metin çıkarılamadı]'
        except Exception as e:
            return f"[DOC okuma hatası: {e}]"
    
    elif ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        slide_texts.append(shape.text_frame.text)
                if slide_texts:
                    text_parts.append(f"--- Slayt {i} ---\n" + '\n'.join(slide_texts))
            return '\n\n'.join(text_parts) if text_parts else '[PPTX dosyasından metin çıkarılamadı]'
        except Exception as e:
            return f"[PPTX okuma hatası: {e}]"
    
    elif ext == '.ppt':
        try:
            import re
            # Use 'strings' command to extract text from binary .ppt file
            result = subprocess.run(
                ['strings', '-e', 'l', file_path],
                capture_output=True, text=True, timeout=10
            )
            text = result.stdout.strip()
            
            if not text:
                result = subprocess.run(
                    ['strings', file_path],
                    capture_output=True, text=True, timeout=10
                )
                text = result.stdout.strip()
            
            if text:
                lines = text.split('\n')
                content_lines = [l.strip() for l in lines if len(l.strip()) > 2
                                 and not re.match(r'^[\x00-\x1f\s]+$', l)]
                return '\n'.join(content_lines) if content_lines else '[PPT dosyasından metin çıkarılamadı]'
            return '[PPT dosyasından metin çıkarılamadı]'
        except Exception as e:
            return f"[PPT okuma hatası: {e}]"
    
    elif ext == '.xlsx':
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else '' for c in row]
                    if any(cells):
                        rows.append('\t'.join(cells))
                if rows:
                    text_parts.append(f"--- Sayfa: {sheet_name} ---\n" + '\n'.join(rows))
            wb.close()
            return '\n\n'.join(text_parts) if text_parts else '[XLSX dosyasından veri çıkarılamadı]'
        except Exception as e:
            return f"[XLSX okuma hatası: {e}]"
    
    elif ext == '.xls':
        try:
            import xlrd
            wb = xlrd.open_workbook(file_path)
            text_parts = []
            for sheet in wb.sheets():
                rows = []
                for row_idx in range(sheet.nrows):
                    cells = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                    if any(c.strip() for c in cells):
                        rows.append('\t'.join(cells))
                if rows:
                    text_parts.append(f"--- Sayfa: {sheet.name} ---\n" + '\n'.join(rows))
            return '\n\n'.join(text_parts) if text_parts else '[XLS dosyasından veri çıkarılamadı]'
        except Exception as e:
            return f"[XLS okuma hatası: {e}]"
    
    else:
        # Universal fallback: try to read as text
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()
            if content.strip():
                return content
            return '[Dosyadan metin çıkarılamadı]'
        except Exception:
            return '[Bu dosya formatı okunamadı]'


@app.route('/api/analyze-file', methods=['POST'])
def analyze_file():
    """Analyze an uploaded document (PDF, DOCX, TXT, etc.)."""
    try:
        if 'file' not in request.files:
            return jsonify({'success': False, 'error': 'Dosya yüklenmedi.'})
        
        file = request.files['file']
        prompt = request.form.get('prompt', 'Bu dosyayı özetle.')
        filename = file.filename or 'dosya'
        ext = os.path.splitext(filename)[1].lower()
        
        # Save temporarily
        import tempfile
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
        file.save(tmp.name)
        tmp.close()
        
        try:
            content = _extract_file_text(tmp.name, ext)
            
            if not content.strip():
                return jsonify({'success': False, 'error': 'Dosyadan metin çıkarılamadı.'})
            
            # Truncate very long files
            if len(content) > 15000:
                content = content[:15000] + '\n\n... (dosya içeriği kısaltıldı)'
            
            # Send to AI with context
            full_prompt = f"Dosya adı: {filename}\n\nDosya içeriği:\n```\n{content}\n```\n\nKullanıcı isteği: {prompt}"
            messages = [{"role": "user", "content": [{"type": "text", "text": full_prompt}]}]
            response = _get_pollinations().generate_response(messages)
            return jsonify({'success': True, 'response': response})
        finally:
            os.unlink(tmp.name)
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)})

@app.route('/api/generate-title', methods=['POST'])
def generate_title():
    """Generate a short chat title using AI based on the conversation start."""
    data = request.get_json()
    user_msg = data.get('user_message', '')
    assistant_msg = data.get('assistant_message', '')

    try:
        prompt = (
            f"Aşağıdaki sohbetin konusunu özetleyen çok kısa bir başlık yaz (en fazla 5 kelime, "
            f"sadece başlığı yaz, başka açıklama yapma, tırnak işareti kullanma):\n"
            f"Kullanıcı: {user_msg[:200]}\n"
            f"Asistan: {assistant_msg[:200]}"
        )
        messages = [{"role": "user", "content": [{"type": "text", "text": prompt}]}]
        title = _get_pollinations().generate_response(messages)
        # Clean up
        title = title.strip().strip('"').strip("'").strip()
        if len(title) > 50:
            title = title[:47] + '...'
        return jsonify({'success': True, 'title': title})
    except Exception as e:
        # Fallback to first few words of user message
        fallback = user_msg[:35] + ('...' if len(user_msg) > 35 else '')
        return jsonify({'success': True, 'title': fallback})


def start_server(port=5789):
    """Start Flask server in a background thread."""
    app.run(host='127.0.0.1', port=port, debug=False, use_reloader=False)

